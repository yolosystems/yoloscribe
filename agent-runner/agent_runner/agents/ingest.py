"""IngestAgent — scheduled agent that processes files in .user/ingest/."""

from __future__ import annotations

import io
import logging
import re
from typing import Callable

from strands import tool
from yoloscribe_io import AgentDefinition, WikiPageMarkdownFile

from .base import BaseAgent
from .search import SearchBackend

log = logging.getLogger(__name__)

_INGEST_PREFIX = ".user/ingest/"
_PROCESSED_PREFIX = ".user/ingest/processed/"

# Matches markdown links pointing at YoloScribe app domains only.
# Used to rewrite absolute YoloScribe URLs to # notation.
# Intentionally narrow — must never match external URLs (foodnetwork.com, etc.).
_YOLOSCRIBE_LINK_RE = re.compile(
    r'\[([^\]]+)\]\(https?://(?:app(?:-dev)?\.yoloscribe\.com|localhost(?::\d+)?)(/[^)]*)\)'
)


def _extract_pdf(file_bytes: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(file_bytes))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(p.strip() for p in pages if p.strip())


def _extract_docx(file_bytes: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    return "\n\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())


def _extract_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [s.text.strip() for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
        if texts:
            slides.append(f"## Slide {i}\n\n" + "\n\n".join(texts))
    return "\n\n".join(slides)


def _extract_xlsx(file_bytes: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    sheets = []
    for name in wb.sheetnames:
        rows = []
        for row in wb[name].iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"## {name}\n\n" + "\n".join(rows))
    return "\n\n".join(sheets)


class IngestAgent(BaseAgent):
    """Agent that processes files queued in .user/ingest/ using search-driven routing.

    Tool surface:
    - ingest_list_pending(): list unprocessed files in .user/ingest/
    - ingest_read(filename): read a specific pending file
    - ingest_mark_processed(filename): move a file to .user/ingest/processed/
    - wiki_search(query): semantic search across the site to find routing targets
    - wiki_list_pages(): list all wiki page paths for structural navigation
    - wiki_read(page_path): read any wiki page for context
    - wiki_write(page_path, content): write to a wiki page
    - notify_owner(message): notify the site owner when content cannot be routed
    - ingest_complete(summary): emit ingest_end event when all files are processed
    - document_index(filename, mode): extract and index an uploaded document
    - http_request + any injected MCP tools
    """

    def __init__(
        self,
        agent_def: AgentDefinition,
        site: str,
        page_path: str,
        storage,
        mcp_tools: list,
        model,
        user_id: str,
        notify_fn: Callable[[str, dict, str], None],
        search: SearchBackend | None = None,
        max_page_reads: int = 10,
        enqueue_fn=None,
    ) -> None:
        super().__init__(
            agent_def=agent_def,
            site=site,
            page_path=page_path,
            storage=storage,
            mcp_tools=mcp_tools,
            model=model,
            user_id=user_id,
            notify_fn=notify_fn,
            search=search,
            max_page_reads=max_page_reads,
        )
        self._read_counter: list[int] = [0]
        self._owner_instructions: str = ""
        self._enqueue_fn = enqueue_fn or (lambda key: None)

    # ── Tool surface ──────────────────────────────────────────────────────────

    @tool
    def ingest_list_pending(self) -> str:
        """List unprocessed files waiting in the ingest queue."""
        prefix = f"{self._site}/{_INGEST_PREFIX}"
        processed_prefix = f"{self._site}/{_PROCESSED_PREFIX}"
        keys = self._storage.list(prefix)
        pending = []
        for key in keys:
            # Exclude the ingest content.md, processed/ subdirectory, and agent files.
            rel = key[len(prefix):]
            if not rel or rel == "content.md" or key.startswith(processed_prefix):
                continue
            if "/.agents/" in rel or rel.startswith(".agents/"):
                continue
            pending.append(rel)
        if not pending:
            return "No pending files."
        return "\n".join(pending)

    @tool
    def ingest_read(self, filename: str) -> str:
        """Read the content of a pending ingest file by its filename."""
        filename = filename.strip().lstrip("/")
        key = f"{self._site}/{_INGEST_PREFIX}{filename}"
        content = self._storage.read(key)
        if content is None:
            return f"File not found: {filename}"
        return content

    @tool
    def ingest_mark_processed(self, filename: str) -> str:
        """Move a processed ingest file to the processed archive."""
        filename = filename.strip().lstrip("/")
        src_key = f"{self._site}/{_INGEST_PREFIX}{filename}"
        dst_key = f"{self._site}/{_PROCESSED_PREFIX}{filename}"
        try:
            self._storage.move(src_key, dst_key)
        except Exception as e:
            return f"Error moving {filename}: {e}"
        log.info("Marked as processed: %s → %s", src_key, dst_key)
        return f"Marked as processed: {filename}"

    @tool
    def wiki_list_pages(self) -> str:
        """List all wiki page paths in this site."""
        prefix = f"{self._site}/"
        keys = self._storage.list(prefix)
        pages = []
        for key in keys:
            if not key.endswith("/content.md"):
                continue
            rel = key[len(prefix):]
            if (
                "/.agents/" in rel
                or "/.user/" in rel
                or "/.archive/" in rel
                or "/.skills/" in rel
                or rel.startswith(".agents/")
                or rel.startswith(".user/")
                or rel.startswith(".archive/")
                or rel.startswith(".skills/")
            ):
                continue
            page_path = rel[: -len("/content.md")] if rel != "content.md" else "(root)"
            pages.append(page_path)
        if not pages:
            return "No wiki pages found."
        return "Wiki pages:\n" + "\n".join(f"- {p}" for p in sorted(pages))

    @tool
    def ingest_complete(self, summary: str) -> str:
        """Signal that all ingest files for this run have been processed.

        Call this once at the end of every run, after all pending files have been
        handled (written, skipped, or flagged for the owner). The summary is
        written to notifications.md so that on_notify agents (e.g. a Discord
        notification agent) can inform the owner.

        Args:
            summary: Plain-text description of what was ingested and where it was
                     routed — e.g. "Processed 2 files: added recipe to cooking/american-south,
                     created new page tech-research-tracker/wasm."
        """
        self._notify("ingest_end", {"summary": summary}, self._user_id)
        log.info("IngestAgent emitted ingest_end: %s", summary[:120])
        return "Ingest complete notification sent."

    @tool
    def notify_owner(self, message: str) -> str:
        """Notify the site owner that a file could not be routed automatically.

        Use this when you cannot determine a suitable destination for an ingest
        file. Do not mark the file as processed after calling this.

        Args:
            message: A plain-text explanation of what was received and why it
                     could not be routed.
        """
        self._notify("ingest_unrouted", {"message": message}, self._user_id)
        log.info("IngestAgent sent ingest_unrouted notification: %s", message[:100])
        return "Owner notified. Leave the file unprocessed so the owner can review it."

    @tool
    def wiki_search(self, query: str) -> str:
        """Search the wiki semantically and return matching page excerpts."""
        results = self._search.search(query, self._site, limit=10)
        if not results:
            return "No matching pages found."
        lines = [
            f"**{r.page_path}** (score: {r.score:.3f})\n{r.excerpt}"
            for r in results
        ]
        return "\n\n".join(lines)

    @tool
    def wiki_read(self, page_path: str) -> str:
        """Read the content of a wiki page by its page path."""
        if self._read_counter[0] >= self._max_page_reads:
            return (
                f"Error: page read limit of {self._max_page_reads} reached. "
                "Complete your task based on what you have already read."
            )
        self._read_counter[0] += 1
        page_path = page_path.strip().strip("/")
        wiki = WikiPageMarkdownFile(site=self._site, page_path=page_path, storage=self._storage)
        return wiki.read()

    @tool
    def wiki_write(self, page_path: str, content: str) -> str:
        """Write content to a wiki page."""
        page_path = page_path.strip().strip("/")
        error = self._check_scope(page_path)
        if error:
            return error
        content = self._rewrite_internal_links(content)
        self._ensure_parent_pages(page_path)
        wiki = WikiPageMarkdownFile(site=self._site, page_path=page_path, storage=self._storage)
        wiki.write(content)
        log.info("IngestAgent wrote to %s/%s", self._site, page_path)
        return f"Written to {page_path}."

    @tool
    def document_index(self, filename: str, mode: str = "both") -> str:
        """Extract text from a binary file in the ingest queue (PDF, DOCX, PPTX, XLSX) and index it.

        Call this instead of ingest_read for binary files. Use the returned text to
        decide where to route the document, then call wiki_write and ingest_mark_processed.

        Args:
            filename: Name of the file in the ingest queue, e.g. 'report.pdf'.
            mode: 'extract' returns text only; 'index-only' indexes without returning text;
                  'both' (default) does both.
        """
        s3_key = f"{self._site}/{_INGEST_PREFIX}{filename}"
        file_bytes = self._storage.read_bytes(s3_key)
        if file_bytes is None:
            return f"File not found: {filename}"

        try:
            extracted_md = self._extract_document(file_bytes, filename)
        except Exception as e:
            return f"Error extracting {filename}: {e}"

        extracted_key = f"{self._site}/{_INGEST_PREFIX}{filename}.extracted.md"
        try:
            self._storage.write_bytes(extracted_key, extracted_md.encode("utf-8"), "text/markdown")
        except Exception as e:
            return f"Error saving extracted text for {filename}: {e}"

        if mode in ("index-only", "both"):
            try:
                self._enqueue_fn(extracted_key)
            except Exception as e:
                log.warning("Failed to enqueue indexing job for %s: %s", extracted_key, e)

        log.info("Indexed document %s → %s (%d chars)", s3_key, extracted_key, len(extracted_md))

        if mode == "index-only":
            return f"Indexed {filename}: {len(extracted_md)} chars extracted, indexing job enqueued."
        return f"Extracted {len(extracted_md)} chars from {filename}.\n\n---\n\n{extracted_md}"

    def _extract_document(self, file_bytes: bytes, filename: str) -> str:
        """Extract text from document bytes, dispatching by file extension."""
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        if ext == "pdf":
            return _extract_pdf(file_bytes)
        if ext == "docx":
            return _extract_docx(file_bytes)
        if ext == "pptx":
            return _extract_pptx(file_bytes)
        if ext in ("xlsx", "xls"):
            return _extract_xlsx(file_bytes)
        return file_bytes.decode("utf-8", errors="replace")

    # ── Content helpers ───────────────────────────────────────────────────────

    def _ensure_parent_pages(self, page_path: str) -> None:
        """Create content.md stubs for any intermediate path segments that lack one.

        e.g. writing to cooking/recipes/heritage-pork will create
        cooking/recipes/content.md if it doesn't exist (cooking/content.md
        is assumed to already exist as the routing target).
        """
        parts = page_path.split("/")
        for i in range(1, len(parts)):
            parent_path = "/".join(parts[:i])
            parent_key = f"{self._site}/{parent_path}/content.md"
            if self._storage.read(parent_key) is None:
                title = parts[i - 1].replace("-", " ").title()
                stub = f"# {title}\n"
                self._storage.write(parent_key, stub)
                log.info("IngestAgent created parent page stub: %s", parent_path)

    def _rewrite_internal_links(self, content: str) -> str:
        """Rewrite absolute YoloScribe app URLs to # notation.

        Only rewrites links pointing at app.yoloscribe.com or app-dev.yoloscribe.com
        (and localhost for dev). External URLs (foodnetwork.com, github.com, etc.)
        are never touched — they must be preserved verbatim.
        """
        site_prefix = f"/{self._site}/"

        def rewrite(m: re.Match) -> str:
            label = m.group(1)
            url_path = m.group(2)  # includes leading slash
            if url_path.startswith(site_prefix):
                url_path = url_path[len(site_prefix):]
            else:
                url_path = url_path.lstrip("/")
            last_seg = url_path.split("/")[-1] if url_path else ""
            if last_seg and "." not in last_seg and re.match(r"^[a-z0-9][a-z0-9/_-]*$", url_path):
                return f"[{label}](#{url_path})"
            return m.group(0)

        return _YOLOSCRIBE_LINK_RE.sub(rewrite, content)

    # ── Scope validation ──────────────────────────────────────────────────────

    def _check_scope(self, page_path: str) -> str | None:
        """Return an error string if page_path is out of scope, else None."""
        if not self.agent_def.scope.matches(page_path):
            return f"Access denied: {page_path} is excluded by this agent's scope settings."
        return None

    # ── Lifecycle ─────────────────────────────────────────────────────────────

    def _read_owner_instructions(self) -> str:
        """Read the owner's routing instructions from .user/ingest/content.md."""
        key = f"{self._site}/.user/ingest/content.md"
        content = self._storage.read(key)
        return (content or "").strip()

    def _build_system_prompt(self) -> str:
        parts = []
        if self.agent_def.description:
            parts.append(self.agent_def.description)
        if self._owner_instructions:
            parts.append(
                "The site owner has provided the following routing instructions. "
                "Follow them precisely — they take priority over your own judgement:\n\n"
                + self._owner_instructions
            )
        parts.append(
            "You process files that arrive in the ingest queue (.user/ingest/) "
            "and route each one to the appropriate wiki page using the wiki's own "
            "structure as the topic guide — no separate index to maintain.\n\n"
            "All content — text and binary documents — flows through the ingest queue. "
            "Binary documents (.pdf, .docx, .pptx, .xlsx) must be extracted before "
            "reading; do not call ingest_read on them.\n\n"
            "Workflow for each run:\n"
            "1. Call wiki_list_pages() to get the full page tree. Always do this first, "
            "before processing any files — it gives you the complete picture of what "
            "topics already exist and prevents creating duplicate or redundant pages.\n"
            "2. Call ingest_list_pending() to find files to process.\n"
            "3. For each pending file:\n"
            "   a. Read the content:\n"
            "      - Binary files (.pdf, .docx, .pptx, .xlsx): call "
            "document_index(filename) to extract text and index it. Use the returned "
            "text for routing decisions. Do NOT call ingest_read on binary files.\n"
            "      - Text files (.md, .txt, .csv, etc.): call ingest_read(filename).\n"
            "   b. Call wiki_search(query) with a concise topic summary to find "
            "semantically similar wiki pages. Higher scores indicate a better match.\n"
            "   c. Choose a destination — strongly prefer routing to an existing page:\n"
            "      - If any existing page (at any depth in the hierarchy) is a "
            "reasonable topical match, append or merge the content there. A partial "
            "match to an existing page is always better than creating a new top-level "
            "topic. When in doubt, add to the closest existing match.\n"
            "      - Only create a new page if you are confident no existing page — "
            "including deep nested pages — is even loosely related to the content.\n"
            "      - Never create a new top-level topic that overlaps with an existing "
            "one at any level of the hierarchy.\n"
            "   d. Call wiki_read(page_path) to read the destination page's current content.\n"
            "   e. Incorporate the new content appropriately (append, merge, or create "
            "a new section), then call wiki_write(page_path, updated_content).\n"
            "   f. If you genuinely cannot determine where the content belongs:\n"
            "      - Call notify_owner(message) describing what you received and why "
            "it could not be routed. Do NOT mark the file as processed.\n"
            "   g. After successfully writing, call ingest_mark_processed(filename).\n"
            "4. When all files have been processed, call ingest_complete(summary) with "
            "a plain-text description of what happened — how many files were processed, "
            "which pages were created or updated, and any files left unprocessed. "
            "Always call this at the end of every run, even if there were no files to process.\n\n"
            "URL hygiene — follow these rules strictly:\n"
            "- Never construct or guess YoloScribe links (app.yoloscribe.com or "
            "app-dev.yoloscribe.com). Internal wiki links use # notation "
            "([label](#page/path)) and must only point to pages you know exist.\n"
            "- External URLs (e.g. foodnetwork.com, github.com, arxiv.org) must be "
            "copied exactly as they appear in the source. Never transform or shorten "
            "them. If the source does not contain a URL, do not add one.\n"
            "- Text extracted from binary files (PDF, DOCX, etc.) may not include "
            "reliable hyperlinks — the extractors capture text only, not the underlying "
            "URLs. Do not invent URLs for titles, authors, or concepts mentioned in "
            "extracted text.\n"
        )
        return "\n\n".join(parts)

    def run(self, prompt: str) -> int:
        self._owner_instructions = self._read_owner_instructions()
        self._notify("ingest_start", {"message": "Processing ingest queue"}, self._user_id)

        tools = [
            self.ingest_list_pending,
            self.ingest_read,
            self.ingest_mark_processed,
            self.wiki_search,
            self.wiki_list_pages,
            self.wiki_read,
            self.wiki_write,
            self.notify_owner,
            self.ingest_complete,
            self.document_index,
        ] + self._mcp_tools

        agent = self._make_strands_agent(tools)
        task = prompt.strip() or "Process all pending ingest files according to your instructions."
        result = agent(task)
        return result.metrics.accumulated_usage.get("totalTokens", 0)
